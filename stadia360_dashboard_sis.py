import streamlit as st
import altair as alt
import pandas as pd
import io
import datetime
from fpdf import FPDF
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from snowflake.snowpark.context import get_active_session

st.set_page_config(
    page_title="STADIA360 統合マーケティングダッシュボード",
    page_icon="📡",
    layout="wide",
)

st.title("STADIA360 統合マーケティングダッシュボード")
st.caption("電通 STADIA360 デモ｜テレビ実視聴データ×多様なKPIデータによる広告効果検証")

session = get_active_session()
DB_SCHEMA = "KFUKAMORI_GEN_DB.STADIA360"


# =====================================================
# データロード
# =====================================================
@st.cache_data(ttl=600)
def load_campaigns():
    return session.sql(f"SELECT * FROM {DB_SCHEMA}.CAMPAIGNS ORDER BY START_DATE").to_pandas()


@st.cache_data(ttl=600)
def load_tv_viewing():
    return session.sql(f"""
        SELECT VIEWING_ID, HOUSEHOLD_ID, CHANNEL, PROGRAM_NAME,
               DATE(VIEWING_START) AS VIEWING_DATE,
               HOUR(VIEWING_START) AS VIEWING_HOUR,
               DAYNAME(VIEWING_START) AS VIEWING_DOW,
               VIEWING_SECONDS, CM_EXPOSED, CAMPAIGN_ID, AREA, DEVICE_TYPE,
               CREATIVE_NAME
        FROM {DB_SCHEMA}.TV_VIEWING_LOG
    """).to_pandas()


@st.cache_data(ttl=600)
def load_cm_cv_joined():
    """CM接触ログとサイトCVを結合したデータ"""
    return session.sql(f"""
        SELECT
            tv.VIEWING_ID,
            tv.HOUSEHOLD_ID,
            tv.CHANNEL,
            tv.PROGRAM_NAME,
            tv.CREATIVE_NAME,
            DATE(tv.VIEWING_START) AS CM_DATE,
            HOUR(tv.VIEWING_START) AS CM_HOUR,
            DAYNAME(tv.VIEWING_START) AS CM_DOW,
            tv.CAMPAIGN_ID,
            tv.AREA,
            cl.CUSTOMER_ID,
            sv.CONVERSION_FLAG,
            DATE(sv.VISIT_TIMESTAMP) AS CV_DATE,
            DATEDIFF('day', DATE(tv.VIEWING_START), DATE(sv.VISIT_TIMESTAMP)) AS DAYS_TO_CV
        FROM {DB_SCHEMA}.TV_VIEWING_LOG tv
        JOIN {DB_SCHEMA}.CUSTOMER_LOYALTY cl ON tv.HOUSEHOLD_ID = cl.HOUSEHOLD_ID
        LEFT JOIN {DB_SCHEMA}.SITE_VISIT_LOG sv ON cl.CUSTOMER_ID = sv.CUSTOMER_ID
        WHERE tv.CM_EXPOSED = TRUE
    """).to_pandas()


@st.cache_data(ttl=600)
def load_loyalty():
    return session.sql(f"SELECT * FROM {DB_SCHEMA}.CUSTOMER_LOYALTY").to_pandas()


@st.cache_data(ttl=600)
def load_attitude():
    return session.sql(f"SELECT * FROM {DB_SCHEMA}.ATTITUDE_CHANGE").to_pandas()


@st.cache_data(ttl=600)
def load_site():
    return session.sql(f"""
        SELECT SESSION_ID, CUSTOMER_ID, DATE(VISIT_TIMESTAMP) AS VISIT_DATE,
               REFERRER_TYPE, REFERRER_DETAIL, PAGE_VIEWS, DURATION_SECONDS,
               CONVERSION_FLAG, CONVERSION_TYPE, CAMPAIGN_ID, DEVICE
        FROM {DB_SCHEMA}.SITE_VISIT_LOG
    """).to_pandas()


@st.cache_data(ttl=600)
def load_purchase():
    return session.sql(f"SELECT * FROM {DB_SCHEMA}.OFFLINE_PURCHASE").to_pandas()


@st.cache_data(ttl=600)
def load_store():
    return session.sql(f"SELECT * FROM {DB_SCHEMA}.STORE_VISIT_LOG").to_pandas()


@st.cache_data(ttl=600)
def load_app_dl():
    return session.sql(f"SELECT * FROM {DB_SCHEMA}.APP_DOWNLOAD_LOG").to_pandas()


@st.cache_data(ttl=600)
def load_app_launch():
    return session.sql(f"""
        SELECT LAUNCH_ID, CUSTOMER_ID, APP_NAME,
               DATE(LAUNCH_TIMESTAMP) AS LAUNCH_DATE,
               SESSION_SECONDS, FEATURES_USED, OS_TYPE, CAMPAIGN_ID
        FROM {DB_SCHEMA}.APP_LAUNCH_LOG
    """).to_pandas()


df_campaigns = load_campaigns()
df_tv = load_tv_viewing()
df_loyalty = load_loyalty()
df_attitude = load_attitude()
df_site = load_site()
df_purchase = load_purchase()
df_store = load_store()
df_app_dl = load_app_dl()
df_app_launch = load_app_launch()
df_cm_cv = load_cm_cv_joined()

# =====================================================
# サイドバーフィルター
# =====================================================
with st.sidebar:
    st.header("フィルター")

    campaign_options = ["全キャンペーン"] + df_campaigns["CAMPAIGN_NAME"].tolist()
    selected_campaign = st.selectbox("キャンペーン", campaign_options)

    area_options = sorted(df_tv["AREA"].unique())
    selected_areas = st.multiselect("エリア", area_options, default=area_options)

    st.markdown("---")
    st.markdown("**データ概要**")
    st.markdown(f"- TV視聴ログ: {len(df_tv):,}件")
    st.markdown(f"- 顧客数: {len(df_loyalty):,}件")
    st.markdown(f"- 態度変容: {len(df_attitude):,}件")
    st.markdown(f"- サイト来訪: {len(df_site):,}件")
    st.markdown(f"- オフライン購買: {len(df_purchase):,}件")
    st.markdown(f"- 来店: {len(df_store):,}件")
    st.markdown(f"- アプリDL: {len(df_app_dl):,}件")
    st.markdown(f"- アプリ起動: {len(df_app_launch):,}件")

# フィルター適用
campaign_id = None
if selected_campaign != "全キャンペーン":
    campaign_id = df_campaigns[df_campaigns["CAMPAIGN_NAME"] == selected_campaign]["CAMPAIGN_ID"].iloc[0]

def filter_by_campaign(df, col="CAMPAIGN_ID"):
    if campaign_id is None:
        return df
    return df[df[col] == campaign_id]

def filter_by_area(df, col="AREA"):
    if col in df.columns:
        return df[df[col].isin(selected_areas)]
    return df

df_tv_f = filter_by_area(filter_by_campaign(df_tv))
df_attitude_f = filter_by_campaign(df_attitude)
df_site_f = filter_by_campaign(df_site)
df_purchase_f = filter_by_area(filter_by_campaign(df_purchase), "STORE_AREA")
df_store_f = filter_by_area(filter_by_campaign(df_store), "STORE_AREA")
df_app_dl_f = filter_by_campaign(df_app_dl)
df_app_launch_f = filter_by_campaign(df_app_launch)
df_cm_cv_f = filter_by_area(filter_by_campaign(df_cm_cv))


# =====================================================
# PDF / Excel レポート生成
# =====================================================
def generate_pdf_report():
    """フィルタ済みデータからPDFレポートを生成"""
    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)

    # --- CJK font setup ---
    # Try to find a CJK-capable font on the system; fall back to DejaVu
    import glob as _glob
    _cjk_font_loaded = False
    _cjk_search_paths = [
        "/usr/share/fonts/**/Noto*CJK*.ttc",
        "/usr/share/fonts/**/Noto*CJK*.ttf",
        "/usr/share/fonts/**/NotoSans*.ttf",
        "/usr/share/fonts/**/DroidSans*.ttf",
        "/usr/share/fonts/**/wqy*.ttc",
        "/usr/share/fonts/**/IPAGothic.ttf",
        "/usr/share/fonts/**/IPAexGothic.ttf",
        "/usr/share/fonts/**/VL-Gothic-Regular.ttf",
        "/usr/share/fonts/**/fonts-japanese-gothic.ttf",
        "/usr/share/fonts/**/*.ttf",
    ]
    for pattern in _cjk_search_paths:
        matches = _glob.glob(pattern, recursive=True)
        if matches:
            try:
                pdf.add_font("cjk", fname=matches[0])
                pdf.add_font("cjk", style="B", fname=matches[0])
                _cjk_font_loaded = True
                break
            except Exception:
                continue

    if not _cjk_font_loaded:
        # Use DejaVu if available (covers Latin/Cyrillic/Greek but not CJK)
        _dejavu_paths = _glob.glob("/usr/share/fonts/**/DejaVuSans.ttf", recursive=True)
        if _dejavu_paths:
            try:
                pdf.add_font("cjk", fname=_dejavu_paths[0])
                pdf.add_font("cjk", style="B", fname=_dejavu_paths[0])
                _cjk_font_loaded = True
            except Exception:
                pass

    # Helper: safe text — encode to latin-1 with replacements if no CJK font
    def safe(text):
        if text is None:
            return ""
        s = str(text)
        if not _cjk_font_loaded:
            # Replace non-latin1 chars with '?'
            return s.encode("latin-1", errors="replace").decode("latin-1")
        return s

    def _set_font(style="", size=10):
        if _cjk_font_loaded:
            pdf.set_font("cjk", style=style.upper(), size=size)
        else:
            pdf.set_font("Helvetica", style, size)

    # Helper: add a section title
    def section_title(title):
        _set_font("B", 14)
        pdf.set_fill_color(41, 128, 185)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 10, safe(title), ln=True, fill=True)
        pdf.set_text_color(0, 0, 0)
        pdf.ln(2)

    # Helper: add a key-value line
    def kv_line(key, value):
        _set_font("B", 10)
        pdf.cell(60, 6, safe(key), ln=False)
        _set_font("", 10)
        pdf.cell(0, 6, safe(value), ln=True)

    # Helper: add a simple table
    def simple_table(headers, rows, col_widths=None):
        if col_widths is None:
            w = 190 / len(headers)
            col_widths = [w] * len(headers)
        # Header
        _set_font("B", 9)
        pdf.set_fill_color(220, 220, 220)
        for i, h in enumerate(headers):
            pdf.cell(col_widths[i], 7, safe(h), border=1, fill=True)
        pdf.ln()
        # Rows
        _set_font("", 9)
        for row in rows[:20]:  # max 20 rows
            for i, val in enumerate(row):
                pdf.cell(col_widths[i], 6, safe(val), border=1)
            pdf.ln()
        if len(rows) > 20:
            _set_font("", 8)
            pdf.cell(0, 6, f"... and {len(rows) - 20} more rows", ln=True)

    today = datetime.date.today().strftime("%Y-%m-%d")

    # ===== Page 1: Cover & KPI =====
    pdf.add_page()
    _set_font("B", 22)
    pdf.cell(0, 20, "STADIA360 Marketing Report", ln=True, align="C")
    _set_font("", 12)
    pdf.cell(0, 10, f"Report Date: {today}", ln=True, align="C")
    pdf.cell(0, 8, safe(f"Campaign: {selected_campaign}"), ln=True, align="C")
    pdf.cell(0, 8, safe(f"Areas: {', '.join(selected_areas)}"), ln=True, align="C")
    pdf.ln(10)

    # KPI Summary
    section_title("KPI Summary")
    cm_count = df_tv_f[df_tv_f["CM_EXPOSED"] == True].shape[0]
    total_tv = df_tv_f.shape[0]
    cm_rate = (cm_count / total_tv * 100) if total_tv > 0 else 0
    cv_count = df_site_f[df_site_f["CONVERSION_FLAG"] == True].shape[0]
    total_sessions = df_site_f.shape[0]
    cvr = (cv_count / total_sessions * 100) if total_sessions > 0 else 0
    store_visits = df_store_f.shape[0]
    avg_stay = df_store_f["STAY_MINUTES"].mean() if store_visits > 0 else 0
    dl_count = df_app_dl_f.shape[0]
    launch_count = df_app_launch_f.shape[0]
    total_purchase = df_purchase_f["AMOUNT"].sum()
    avg_nps = df_loyalty["NPS_SCORE"].mean()
    total_pv = df_site_f["PAGE_VIEWS"].sum()
    total_hours = df_tv_f["VIEWING_SECONDS"].sum() / 3600

    kv_line("CM Contacts:", f"{cm_count:,} (Rate: {cm_rate:.1f}%)")
    kv_line("Site CV:", f"{cv_count:,} (CVR: {cvr:.1f}%)")
    kv_line("Store Visits:", f"{store_visits:,} (Avg Stay: {avg_stay:.0f} min)")
    kv_line("App DL / Launch:", f"{dl_count:,} / {launch_count:,}")
    kv_line("Purchase Total:", f"JPY {total_purchase:,.0f}")
    kv_line("Avg NPS:", f"{avg_nps:.1f}")
    kv_line("Total PV:", f"{total_pv:,}")
    kv_line("Total View Hours:", f"{total_hours:,.0f} hrs")
    pdf.ln(5)

    # Campaign List
    section_title("Campaigns")
    camp_rows = []
    for _, r in df_campaigns.iterrows():
        camp_rows.append([
            safe(r["CAMPAIGN_ID"]), safe(r["CAMPAIGN_NAME"]),
            safe(r["START_DATE"]), safe(r["END_DATE"]),
            f"{r['BUDGET_MM']}M JPY",
        ])
    simple_table(
        ["ID", "Name", "Start", "End", "Budget"],
        camp_rows,
        [20, 70, 30, 30, 40],
    )
    pdf.ln(5)

    # ===== Page 2: TV & CM Analysis =====
    pdf.add_page()
    section_title("TV Viewing by Channel")
    ch_data = df_tv_f.groupby("CHANNEL").agg(
        VIEWS=("VIEWING_ID", "count"), CM=("CM_EXPOSED", "sum")
    ).reset_index()
    ch_data["Rate"] = (ch_data["CM"] / ch_data["VIEWS"] * 100).round(1)
    ch_rows = [[safe(r["CHANNEL"]), f"{r['VIEWS']:,}", f"{r['CM']:,}", f"{r['Rate']:.1f}%"]
               for _, r in ch_data.sort_values("VIEWS", ascending=False).iterrows()]
    simple_table(["Channel", "Views", "CM Contacts", "CM Rate"], ch_rows, [50, 45, 45, 50])
    pdf.ln(5)

    # CM Effect
    section_title("CM Effect Analysis (CV Rate)")
    df_cm_cv_tmp = df_cm_cv_f.copy()
    df_cm_cv_tmp["CV"] = pd.to_numeric(df_cm_cv_tmp["CONVERSION_FLAG"], errors="coerce").fillna(0).astype(int)
    ch_cv = df_cm_cv_tmp.groupby("CHANNEL").agg(
        Contacts=("VIEWING_ID", "count"), CVs=("CV", "sum")
    ).reset_index()
    ch_cv["CVR"] = (ch_cv["CVs"] / ch_cv["Contacts"] * 100).round(2)
    cv_rows = [[safe(r["CHANNEL"]), f"{r['Contacts']:,}", f"{r['CVs']:,}", f"{r['CVR']:.2f}%"]
               for _, r in ch_cv.sort_values("CVR", ascending=False).iterrows()]
    simple_table(["Channel", "CM Contacts", "CVs", "CV Rate"], cv_rows, [50, 45, 45, 50])
    pdf.ln(5)

    # Creative drill
    section_title("Creative CV Rate")
    cr_cv = df_cm_cv_tmp[df_cm_cv_tmp["CREATIVE_NAME"].notna()].groupby("CREATIVE_NAME").agg(
        Contacts=("VIEWING_ID", "count"), CVs=("CV", "sum")
    ).reset_index()
    cr_cv["CVR"] = (cr_cv["CVs"] / cr_cv["Contacts"] * 100).round(2)
    cr_rows = [[safe(r["CREATIVE_NAME"]), f"{r['Contacts']:,}", f"{r['CVs']:,}", f"{r['CVR']:.2f}%"]
               for _, r in cr_cv.sort_values("CVR", ascending=False).iterrows()]
    simple_table(["Creative", "Contacts", "CVs", "CV Rate"], cr_rows, [60, 40, 40, 50])
    pdf.ln(5)

    # Program Ranking
    section_title("Program Ranking (Top 10 by CV Rate)")
    pgm = df_cm_cv_tmp.groupby("PROGRAM_NAME").agg(
        Contacts=("VIEWING_ID", "count"), CVs=("CV", "sum")
    ).reset_index()
    pgm["CVR"] = (pgm["CVs"] / pgm["Contacts"] * 100).round(2)
    pgm = pgm.sort_values("CVR", ascending=False).head(10)
    pgm_rows = [[safe(r["PROGRAM_NAME"]), f"{r['Contacts']:,}", f"{r['CVs']:,}", f"{r['CVR']:.2f}%"]
                for _, r in pgm.iterrows()]
    simple_table(["Program", "Contacts", "CVs", "CV Rate"], pgm_rows, [60, 40, 40, 50])

    # ===== Page 3: Attitude, Site, Purchase =====
    pdf.add_page()
    section_title("Attitude Change (Lift)")
    exposed = df_attitude_f[df_attitude_f["CM_EXPOSED"] == True]
    unexposed = df_attitude_f[df_attitude_f["CM_EXPOSED"] == False]
    lift_rows = []
    for stage, bef, aft in [("Awareness", "AWARENESS_BEFORE", "AWARENESS_AFTER"),
                             ("Interest", "INTEREST_BEFORE", "INTEREST_AFTER"),
                             ("Consider", "CONSIDER_BEFORE", "CONSIDER_AFTER"),
                             ("Purchase", "PURCHASE_BEFORE", "PURCHASE_AFTER")]:
        e_lift = (exposed[aft] - exposed[bef]).mean() if len(exposed) > 0 else 0
        u_lift = (unexposed[aft] - unexposed[bef]).mean() if len(unexposed) > 0 else 0
        lift_rows.append([stage, f"{e_lift:.2f} pt", f"{u_lift:.2f} pt", f"{e_lift - u_lift:.2f} pt"])
    simple_table(["Stage", "CM Exposed", "Unexposed", "Diff"], lift_rows, [40, 50, 50, 50])
    pdf.ln(5)

    # Site Visit
    section_title("Site Visit Summary")
    ref_data = df_site_f.groupby("REFERRER_TYPE").agg(
        Sessions=("SESSION_ID", "count"), CVs=("CONVERSION_FLAG", "sum"), AvgPV=("PAGE_VIEWS", "mean")
    ).reset_index()
    ref_data["CVR"] = (ref_data["CVs"] / ref_data["Sessions"] * 100).round(1)
    ref_rows = [[safe(r["REFERRER_TYPE"]), f"{r['Sessions']:,}", f"{r['CVs']:,}",
                 f"{r['CVR']:.1f}%", f"{r['AvgPV']:.1f}"]
                for _, r in ref_data.sort_values("Sessions", ascending=False).iterrows()]
    simple_table(["Referrer", "Sessions", "CVs", "CVR", "Avg PV"], ref_rows, [40, 35, 30, 35, 50])
    pdf.ln(5)

    # Purchase
    section_title("Offline Purchase Summary")
    cat_data = df_purchase_f.groupby("PRODUCT_CATEGORY").agg(
        Total=("AMOUNT", "sum"), Count=("PURCHASE_ID", "count")
    ).reset_index()
    cat_data["Avg"] = (cat_data["Total"] / cat_data["Count"]).round(0)
    pur_rows = [[safe(r["PRODUCT_CATEGORY"]), f"JPY {r['Total']:,.0f}", f"{r['Count']:,}", f"JPY {r['Avg']:,.0f}"]
                for _, r in cat_data.sort_values("Total", ascending=False).iterrows()]
    simple_table(["Category", "Total Amount", "Count", "Avg Amount"], pur_rows, [50, 50, 40, 50])
    pdf.ln(3)
    cm_exp_avg = df_purchase_f[df_purchase_f["CM_EXPOSED"] == True]["AMOUNT"].mean() if len(df_purchase_f[df_purchase_f["CM_EXPOSED"] == True]) > 0 else 0
    cm_unexp_avg = df_purchase_f[df_purchase_f["CM_EXPOSED"] == False]["AMOUNT"].mean() if len(df_purchase_f[df_purchase_f["CM_EXPOSED"] == False]) > 0 else 0
    kv_line("CM Exposed Avg:", f"JPY {cm_exp_avg:,.0f}")
    kv_line("Unexposed Avg:", f"JPY {cm_unexp_avg:,.0f}")

    # ===== Page 4: Store, App, Loyalty =====
    pdf.add_page()
    section_title("Store Visits")
    store_data = df_store_f.groupby("STORE_NAME").agg(
        Visits=("VISIT_ID", "count"), AvgStay=("STAY_MINUTES", "mean")
    ).reset_index()
    store_data["AvgStay"] = store_data["AvgStay"].round(1)
    st_rows = [[safe(r["STORE_NAME"]), f"{r['Visits']:,}", f"{r['AvgStay']:.1f} min"]
               for _, r in store_data.sort_values("Visits", ascending=False).iterrows()]
    simple_table(["Store", "Visits", "Avg Stay"], st_rows, [80, 55, 55])
    pdf.ln(5)

    section_title("App Downloads & Launches")
    dl_by_app = df_app_dl_f.groupby("APP_NAME").agg(DL=("DOWNLOAD_ID", "count")).reset_index()
    launch_by_app = df_app_launch_f.groupby("APP_NAME").agg(Launch=("LAUNCH_ID", "count")).reset_index()
    app_data = dl_by_app.merge(launch_by_app, on="APP_NAME", how="outer").fillna(0)
    app_rows = [[safe(r["APP_NAME"]), f"{int(r['DL']):,}", f"{int(r['Launch']):,}"]
                for _, r in app_data.iterrows()]
    simple_table(["App", "Downloads", "Launches"], app_rows, [80, 55, 55])
    pdf.ln(5)

    section_title("Customer Loyalty")
    seg_data = df_loyalty.groupby("LOYALTY_SEGMENT").agg(
        Count=("CUSTOMER_ID", "count"), AvgNPS=("NPS_SCORE", "mean"), AvgLTV=("LTV_AMOUNT", "mean")
    ).reset_index()
    seg_rows = [[safe(r["LOYALTY_SEGMENT"]), f"{r['Count']:,}", f"{r['AvgNPS']:.1f}", f"JPY {r['AvgLTV']:,.0f}"]
                for _, r in seg_data.iterrows()]
    simple_table(["Segment", "Count", "Avg NPS", "Avg LTV"], seg_rows, [50, 40, 50, 50])

    # Output — convert bytearray to bytes for st.download_button compatibility
    return bytes(pdf.output())


def generate_excel_report():
    """フィルタ済みデータからExcelレポートを生成"""
    output = io.BytesIO()
    wb = Workbook()

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="2980B9", end_color="2980B9", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    def write_sheet(ws, df, sheet_name=None):
        if sheet_name:
            ws.title = sheet_name
        # Header
        for col_idx, col_name in enumerate(df.columns, 1):
            cell = ws.cell(row=1, column=col_idx, value=str(col_name))
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            cell.border = thin_border
        # Data
        for row_idx, row in enumerate(df.itertuples(index=False), 2):
            for col_idx, val in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=val)
                cell.border = thin_border
        # Auto-width
        for col_idx, col_name in enumerate(df.columns, 1):
            max_len = max(len(str(col_name)), df[col_name].astype(str).str.len().max() if len(df) > 0 else 0)
            ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = min(max_len + 4, 40)

    # Sheet 1: KPI
    ws_kpi = wb.active
    ws_kpi.title = "KPI Summary"
    cm_count = df_tv_f[df_tv_f["CM_EXPOSED"] == True].shape[0]
    total_tv = df_tv_f.shape[0]
    cv_count = df_site_f[df_site_f["CONVERSION_FLAG"] == True].shape[0]
    total_sessions = df_site_f.shape[0]
    kpi_data = pd.DataFrame({
        "KPI": ["CM Contacts", "CM Rate(%)", "Site CV", "CVR(%)",
                "Store Visits", "App DL", "App Launch",
                "Purchase Total(JPY)", "Avg NPS", "Total PV"],
        "Value": [
            cm_count, round(cm_count / total_tv * 100, 1) if total_tv > 0 else 0,
            cv_count, round(cv_count / total_sessions * 100, 1) if total_sessions > 0 else 0,
            df_store_f.shape[0], df_app_dl_f.shape[0], df_app_launch_f.shape[0],
            df_purchase_f["AMOUNT"].sum(), round(df_loyalty["NPS_SCORE"].mean(), 1),
            df_site_f["PAGE_VIEWS"].sum(),
        ],
    })
    write_sheet(ws_kpi, kpi_data, "KPI Summary")

    # Sheet 2: Campaigns
    ws_camp = wb.create_sheet("Campaigns")
    write_sheet(ws_camp, df_campaigns[["CAMPAIGN_ID", "CAMPAIGN_NAME", "ADVERTISER",
                                       "START_DATE", "END_DATE", "BUDGET_MM", "PRODUCT_CATEGORY",
                                       "TARGET_AREA"]], "Campaigns")

    # Sheet 3: Channel CM Effect
    ws_ch = wb.create_sheet("Channel CM Effect")
    df_cm_tmp = df_cm_cv_f.copy()
    df_cm_tmp["CV"] = pd.to_numeric(df_cm_tmp["CONVERSION_FLAG"], errors="coerce").fillna(0).astype(int)
    ch_effect = df_cm_tmp.groupby("CHANNEL").agg(
        CM_Contacts=("VIEWING_ID", "count"), CVs=("CV", "sum")
    ).reset_index()
    ch_effect["CV_Rate(%)"] = (ch_effect["CVs"] / ch_effect["CM_Contacts"] * 100).round(2)
    write_sheet(ws_ch, ch_effect.sort_values("CV_Rate(%)", ascending=False), "Channel CM Effect")

    # Sheet 4: Creative CV
    ws_cr = wb.create_sheet("Creative CV")
    cr_effect = df_cm_tmp[df_cm_tmp["CREATIVE_NAME"].notna()].groupby("CREATIVE_NAME").agg(
        CM_Contacts=("VIEWING_ID", "count"), CVs=("CV", "sum")
    ).reset_index()
    cr_effect["CV_Rate(%)"] = (cr_effect["CVs"] / cr_effect["CM_Contacts"] * 100).round(2)
    write_sheet(ws_cr, cr_effect.sort_values("CV_Rate(%)", ascending=False), "Creative CV")

    # Sheet 5: Site Visit
    ws_site = wb.create_sheet("Site Visit")
    ref_effect = df_site_f.groupby("REFERRER_TYPE").agg(
        Sessions=("SESSION_ID", "count"), CVs=("CONVERSION_FLAG", "sum"),
        Avg_PV=("PAGE_VIEWS", "mean")
    ).reset_index()
    ref_effect["CVR(%)"] = (ref_effect["CVs"] / ref_effect["Sessions"] * 100).round(1)
    ref_effect["Avg_PV"] = ref_effect["Avg_PV"].round(1)
    write_sheet(ws_site, ref_effect.sort_values("Sessions", ascending=False), "Site Visit")

    # Sheet 6: Purchase
    ws_pur = wb.create_sheet("Purchase")
    pur_effect = df_purchase_f.groupby("PRODUCT_CATEGORY").agg(
        Total_Amount=("AMOUNT", "sum"), Count=("PURCHASE_ID", "count")
    ).reset_index()
    pur_effect["Avg_Amount"] = (pur_effect["Total_Amount"] / pur_effect["Count"]).round(0)
    write_sheet(ws_pur, pur_effect.sort_values("Total_Amount", ascending=False), "Purchase")

    # Sheet 7: Loyalty
    ws_loy = wb.create_sheet("Loyalty")
    loy_data = df_loyalty.groupby("LOYALTY_SEGMENT").agg(
        Count=("CUSTOMER_ID", "count"), Avg_NPS=("NPS_SCORE", "mean"), Avg_LTV=("LTV_AMOUNT", "mean")
    ).reset_index()
    loy_data["Avg_NPS"] = loy_data["Avg_NPS"].round(1)
    loy_data["Avg_LTV"] = loy_data["Avg_LTV"].round(0)
    write_sheet(ws_loy, loy_data, "Loyalty")

    wb.save(output)
    output.seek(0)
    return output.getvalue()


def generate_pptx_report():
    """フィルタ済みデータからPPTXレポートを生成（matplotlibグラフ付き）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(41, 128, 185)
    DARK = RGBColor(44, 62, 80)
    WHITE = RGBColor(255, 255, 255)
    GRAY_BG = RGBColor(240, 240, 240)
    # matplotlib color palette
    MPL_BLUE = "#2980B9"
    MPL_ORANGE = "#E67E22"
    MPL_GREEN = "#27AE60"
    MPL_RED = "#E74C3C"
    MPL_PURPLE = "#8E44AD"
    MPL_COLORS = [MPL_BLUE, MPL_ORANGE, MPL_GREEN, MPL_RED, MPL_PURPLE, "#3498DB", "#F39C12", "#1ABC9C"]

    def _add_slide():
        layout = prs.slide_layouts[6]  # Blank
        return prs.slides.add_slide(layout)

    def _add_title_bar(slide, text, top=Inches(0.3)):
        txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BLUE
        return top + Inches(0.7)

    def _fig_to_image(fig):
        """matplotlibのfigをBytesIOのPNG画像に変換"""
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf

    def _add_chart_image(slide, fig, left=Inches(0.5), top=Inches(1.1), width=Inches(6), height=Inches(5.5)):
        """matplotlibのfigをスライドに画像として挿入"""
        img_buf = _fig_to_image(fig)
        slide.shapes.add_picture(img_buf, left, top, width, height)

    def _add_table(slide, headers, rows, left=Inches(0.5), top=Inches(1.2), width=Inches(5.8)):
        n_rows = min(len(rows), 20) + 1
        n_cols = len(headers)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.3 * n_rows))
        tbl = tbl_shape.table
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = str(h)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        for r_idx, row in enumerate(rows[:20]):
            for c_idx, val in enumerate(row):
                cell = tbl.cell(r_idx + 1, c_idx)
                cell.text = str(val) if val is not None else ""
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(8)
                    paragraph.alignment = PP_ALIGN.CENTER
                if r_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GRAY_BG
        return tbl_shape

    def _add_kv_box(slide, items, left=Inches(0.5), top=Inches(1.2)):
        txBox = slide.shapes.add_textbox(left, top, Inches(5.8), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (k, v) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run_k = p.add_run()
            run_k.text = f"{k}  "
            run_k.font.size = Pt(16)
            run_k.font.bold = True
            run_k.font.color.rgb = DARK
            run_v = p.add_run()
            run_v.text = str(v)
            run_v.font.size = Pt(16)
            run_v.font.color.rgb = BLUE
            p.space_after = Pt(8)

    today = datetime.date.today().strftime("%Y-%m-%d")

    # ===== Slide 1: Title =====
    slide = _add_slide()
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(23, 37, 84)
    bg.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "STADIA360"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Marketing Analytics Report"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(147, 197, 253)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = f"\nReport Date: {today}"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(200, 200, 200)
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph()
    p4.text = f"Campaign: {selected_campaign}"
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(200, 200, 200)
    p4.alignment = PP_ALIGN.CENTER
    p5 = tf.add_paragraph()
    p5.text = f"Areas: {', '.join(selected_areas)}"
    p5.font.size = Pt(16)
    p5.font.color.rgb = RGBColor(200, 200, 200)
    p5.alignment = PP_ALIGN.CENTER

    # ===== Slide 2: KPI Summary (KV left + bar chart right) =====
    slide = _add_slide()
    _add_title_bar(slide, "KPI Summary")
    cm_count = df_tv_f[df_tv_f["CM_EXPOSED"] == True].shape[0]
    total_tv = df_tv_f.shape[0]
    cm_rate = (cm_count / total_tv * 100) if total_tv > 0 else 0
    cv_count = df_site_f[df_site_f["CONVERSION_FLAG"] == True].shape[0]
    total_sessions = df_site_f.shape[0]
    cvr = (cv_count / total_sessions * 100) if total_sessions > 0 else 0
    store_visits = df_store_f.shape[0]
    avg_stay = df_store_f["STAY_MINUTES"].mean() if store_visits > 0 else 0
    dl_count = df_app_dl_f.shape[0]
    launch_count = df_app_launch_f.shape[0]
    total_purchase = df_purchase_f["AMOUNT"].sum()
    avg_nps = df_loyalty["NPS_SCORE"].mean()
    total_pv = df_site_f["PAGE_VIEWS"].sum()
    total_hours = df_tv_f["VIEWING_SECONDS"].sum() / 3600
    _add_kv_box(slide, [
        ("CM Contacts:", f"{cm_count:,} ({cm_rate:.1f}%)"),
        ("Site CV:", f"{cv_count:,} (CVR: {cvr:.1f}%)"),
        ("Store Visits:", f"{store_visits:,} ({avg_stay:.0f} min avg)"),
        ("App DL / Launch:", f"{dl_count:,} / {launch_count:,}"),
        ("Purchase Total:", f"JPY {total_purchase:,.0f}"),
        ("Avg NPS:", f"{avg_nps:.1f}"),
        ("Total PV:", f"{total_pv:,}"),
        ("View Hours:", f"{total_hours:,.0f} hrs"),
    ])
    # KPI bar chart on right
    fig, ax = plt.subplots(figsize=(5, 4))
    kpi_labels = ["CM Rate%", "CVR%", "Avg NPS"]
    kpi_vals = [cm_rate, cvr, avg_nps]
    bars = ax.bar(kpi_labels, kpi_vals, color=[MPL_BLUE, MPL_ORANGE, MPL_GREEN])
    ax.set_ylabel("Value")
    ax.set_title("Key Metrics Overview")
    for bar, val in zip(bars, kpi_vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5, f"{val:.1f}", ha="center", fontsize=10)
    fig.tight_layout()
    _add_chart_image(slide, fig, left=Inches(7), top=Inches(1.2), width=Inches(5.8), height=Inches(5))

    # ===== Slide 3: TV Viewing by Channel (bar chart + table) =====
    slide = _add_slide()
    _add_title_bar(slide, "TV Viewing by Channel")
    ch_data = df_tv_f.groupby("CHANNEL").agg(
        VIEWS=("VIEWING_ID", "count"), CM=("CM_EXPOSED", "sum")
    ).reset_index()
    ch_data["Rate"] = (ch_data["CM"] / ch_data["VIEWS"] * 100).round(1)
    ch_data = ch_data.sort_values("VIEWS", ascending=True)
    # Chart
    fig, ax = plt.subplots(figsize=(5.5, 4.5))
    y_pos = range(len(ch_data))
    ax.barh(list(y_pos), ch_data["VIEWS"].values, color=MPL_BLUE, alpha=0.7, label="Total Views")
    ax.barh(list(y_pos), ch_data["CM"].values, color=MPL_ORANGE, alpha=0.9, label="CM Contacts")
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(ch_data["CHANNEL"].values, fontsize=9)
    ax.set_xlabel("Count")
    ax.set_title("Views & CM Contacts by Channel")
    ax.legend(loc="lower right", fontsize=8)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    fig.tight_layout()
    _add_chart_image(slide, fig, left=Inches(0.5), top=Inches(1.1), width=Inches(6.3), height=Inches(5.5))
    # Table on right
    ch_rows = [[r["CHANNEL"], f"{r['VIEWS']:,}", f"{r['CM']:,}", f"{r['Rate']:.1f}%"]
               for _, r in ch_data.sort_values("VIEWS", ascending=False).iterrows()]
    _add_table(slide, ["Channel", "Views", "CM", "Rate"], ch_rows, left=Inches(7.2), top=Inches(1.2), width=Inches(5.6))

    # ===== Slide 4: CM Effect Analysis (grouped bar + creative table) =====
    slide = _add_slide()
    _add_title_bar(slide, "CM Effect Analysis (CV Rate)")
    df_cm_cv_tmp = df_cm_cv_f.copy()
    df_cm_cv_tmp["CV"] = pd.to_numeric(df_cm_cv_tmp["CONVERSION_FLAG"], errors="coerce").fillna(0).astype(int)
    ch_cv = df_cm_cv_tmp.groupby("CHANNEL").agg(
        Contacts=("VIEWING_ID", "count"), CVs=("CV", "sum")
    ).reset_index()
    ch_cv["CVR"] = (ch_cv["CVs"] / ch_cv["Contacts"] * 100).round(2)
    ch_cv = ch_cv.sort_values("CVR", ascending=False)
    # Chart
    fig, ax = plt.subplots(figsize=(5.5, 4.5))
    x = range(len(ch_cv))
    ax.bar(list(x), ch_cv["Contacts"].values, color=MPL_BLUE, alpha=0.7, label="CM Contacts")
    ax2 = ax.twinx()
    ax2.plot(list(x), ch_cv["CVR"].values, color=MPL_RED, marker="o", linewidth=2, label="CV Rate %")
    ax.set_xticks(list(x))
    ax.set_xticklabels(ch_cv["CHANNEL"].values, fontsize=8, rotation=30, ha="right")
    ax.set_ylabel("CM Contacts")
    ax2.set_ylabel("CV Rate (%)")
    ax.set_title("CM Contacts & CV Rate by Channel")
    ax.legend(loc="upper left", fontsize=8)
    ax2.legend(loc="upper right", fontsize=8)
    fig.tight_layout()
    _add_chart_image(slide, fig, left=Inches(0.5), top=Inches(1.1), width=Inches(6.3), height=Inches(5.5))
    # Creative table on right
    cr_cv = df_cm_cv_tmp[df_cm_cv_tmp["CREATIVE_NAME"].notna()].groupby("CREATIVE_NAME").agg(
        Contacts=("VIEWING_ID", "count"), CVs=("CV", "sum")
    ).reset_index()
    cr_cv["CVR"] = (cr_cv["CVs"] / cr_cv["Contacts"] * 100).round(2)
    cr_rows = [[r["CREATIVE_NAME"], f"{r['Contacts']:,}", f"{r['CVs']:,}", f"{r['CVR']:.2f}%"]
               for _, r in cr_cv.sort_values("CVR", ascending=False).iterrows()]
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(1.0), Inches(5), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Creative CV Rate"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE
    _add_table(slide, ["Creative", "Contacts", "CVs", "CVR"], cr_rows, left=Inches(7.2), top=Inches(1.4), width=Inches(5.6))

    # ===== Slide 5: Attitude Change (grouped bar chart) =====
    slide = _add_slide()
    _add_title_bar(slide, "Attitude Change (CM Exposed vs Unexposed)")
    exposed = df_attitude_f[df_attitude_f["CM_EXPOSED"] == True]
    unexposed = df_attitude_f[df_attitude_f["CM_EXPOSED"] == False]
    stages = ["Awareness", "Interest", "Consider", "Purchase"]
    cols_bef_aft = [("AWARENESS_BEFORE", "AWARENESS_AFTER"),
                    ("INTEREST_BEFORE", "INTEREST_AFTER"),
                    ("CONSIDER_BEFORE", "CONSIDER_AFTER"),
                    ("PURCHASE_BEFORE", "PURCHASE_AFTER")]
    e_lifts, u_lifts = [], []
    for bef, aft in cols_bef_aft:
        e_lifts.append((exposed[aft] - exposed[bef]).mean() if len(exposed) > 0 else 0)
        u_lifts.append((unexposed[aft] - unexposed[bef]).mean() if len(unexposed) > 0 else 0)
    # Chart
    fig, ax = plt.subplots(figsize=(6, 4.5))
    x = range(len(stages))
    w = 0.35
    ax.bar([i - w / 2 for i in x], e_lifts, w, label="CM Exposed", color=MPL_BLUE)
    ax.bar([i + w / 2 for i in x], u_lifts, w, label="Unexposed", color=MPL_ORANGE)
    ax.set_xticks(list(x))
    ax.set_xticklabels(stages, fontsize=10)
    ax.set_ylabel("Lift (pt)")
    ax.set_title("Attitude Change Lift: Exposed vs Unexposed")
    ax.legend(fontsize=9)
    for i in x:
        ax.text(i - w / 2, e_lifts[i] + 0.2, f"{e_lifts[i]:.1f}", ha="center", fontsize=8)
        ax.text(i + w / 2, u_lifts[i] + 0.2, f"{u_lifts[i]:.1f}", ha="center", fontsize=8)
    fig.tight_layout()
    _add_chart_image(slide, fig, left=Inches(0.5), top=Inches(1.1), width=Inches(6.3), height=Inches(5.5))
    # Table on right
    lift_rows = []
    for i, stage in enumerate(stages):
        diff = e_lifts[i] - u_lifts[i]
        lift_rows.append([stage, f"{e_lifts[i]:.2f}", f"{u_lifts[i]:.2f}", f"{diff:.2f}"])
    _add_table(slide, ["Stage", "Exposed", "Unexposed", "Diff"], lift_rows, left=Inches(7.2), top=Inches(1.2), width=Inches(5.6))

    # ===== Slide 6: Site Visit + Purchase =====
    slide = _add_slide()
    _add_title_bar(slide, "Site Visit & Offline Purchase")
    ref_data = df_site_f.groupby("REFERRER_TYPE").agg(
        Sessions=("SESSION_ID", "count"), CVs=("CONVERSION_FLAG", "sum"), AvgPV=("PAGE_VIEWS", "mean")
    ).reset_index()
    ref_data["CVR"] = (ref_data["CVs"] / ref_data["Sessions"] * 100).round(1)
    ref_data = ref_data.sort_values("Sessions", ascending=False)
    # Site visit pie chart
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    ax1.pie(ref_data["Sessions"].values, labels=ref_data["REFERRER_TYPE"].values,
            autopct="%1.1f%%", colors=MPL_COLORS[:len(ref_data)], startangle=90)
    ax1.set_title("Site Sessions by Referrer", fontsize=11)
    # Purchase bar chart
    cat_data = df_purchase_f.groupby("PRODUCT_CATEGORY").agg(
        Total=("AMOUNT", "sum"), Count=("PURCHASE_ID", "count")
    ).reset_index().sort_values("Total", ascending=True)
    ax2.barh(cat_data["PRODUCT_CATEGORY"].values, cat_data["Total"].values, color=MPL_GREEN)
    ax2.set_xlabel("Total Amount (JPY)")
    ax2.set_title("Purchase by Category", fontsize=11)
    ax2.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x / 1e6:.1f}M"))
    fig.tight_layout()
    _add_chart_image(slide, fig, left=Inches(0.3), top=Inches(1.1), width=Inches(12.5), height=Inches(5.5))

    # ===== Slide 7: Store Visit + App (charts side by side) =====
    slide = _add_slide()
    _add_title_bar(slide, "Store Visit & App Analysis")
    store_data = df_store_f.groupby("STORE_NAME").agg(
        Visits=("VISIT_ID", "count"), AvgStay=("STAY_MINUTES", "mean")
    ).reset_index().sort_values("Visits", ascending=True)
    dl_by_app = df_app_dl_f.groupby("APP_NAME").agg(DL=("DOWNLOAD_ID", "count")).reset_index()
    launch_by_app = df_app_launch_f.groupby("APP_NAME").agg(Launch=("LAUNCH_ID", "count")).reset_index()
    app_data = dl_by_app.merge(launch_by_app, on="APP_NAME", how="outer").fillna(0)
    # Charts
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    # Store visits
    ax1.barh(store_data["STORE_NAME"].values, store_data["Visits"].values, color=MPL_BLUE)
    ax1.set_xlabel("Visit Count")
    ax1.set_title("Store Visits", fontsize=11)
    # App DL vs Launch
    x = range(len(app_data))
    w = 0.35
    ax2.bar([i - w / 2 for i in x], app_data["DL"].values.astype(int), w, label="Downloads", color=MPL_ORANGE)
    ax2.bar([i + w / 2 for i in x], app_data["Launch"].values.astype(int), w, label="Launches", color=MPL_PURPLE)
    ax2.set_xticks(list(x))
    ax2.set_xticklabels(app_data["APP_NAME"].values, fontsize=8, rotation=30, ha="right")
    ax2.set_ylabel("Count")
    ax2.set_title("App Downloads & Launches", fontsize=11)
    ax2.legend(fontsize=8)
    fig.tight_layout()
    _add_chart_image(slide, fig, left=Inches(0.3), top=Inches(1.1), width=Inches(12.5), height=Inches(5.5))

    # ===== Slide 8: Loyalty (pie + bar) =====
    slide = _add_slide()
    _add_title_bar(slide, "Customer Loyalty")
    seg_data = df_loyalty.groupby("LOYALTY_SEGMENT").agg(
        Count=("CUSTOMER_ID", "count"), AvgNPS=("NPS_SCORE", "mean"), AvgLTV=("LTV_AMOUNT", "mean")
    ).reset_index()
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    # Segment distribution pie
    ax1.pie(seg_data["Count"].values, labels=seg_data["LOYALTY_SEGMENT"].values,
            autopct="%1.1f%%", colors=MPL_COLORS[:len(seg_data)], startangle=90)
    ax1.set_title("Loyalty Segment Distribution", fontsize=11)
    # NPS & LTV bar
    x = range(len(seg_data))
    ax2.bar(list(x), seg_data["AvgNPS"].values, color=MPL_BLUE, alpha=0.8, label="Avg NPS")
    ax2b = ax2.twinx()
    ax2b.plot(list(x), seg_data["AvgLTV"].values, color=MPL_RED, marker="s", linewidth=2, label="Avg LTV")
    ax2.set_xticks(list(x))
    ax2.set_xticklabels(seg_data["LOYALTY_SEGMENT"].values, fontsize=9)
    ax2.set_ylabel("Avg NPS")
    ax2b.set_ylabel("Avg LTV (JPY)")
    ax2.set_title("NPS & LTV by Segment", fontsize=11)
    ax2.legend(loc="upper left", fontsize=8)
    ax2b.legend(loc="upper right", fontsize=8)
    ax2b.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x / 1e3:.0f}K"))
    fig.tight_layout()
    _add_chart_image(slide, fig, left=Inches(0.3), top=Inches(1.1), width=Inches(12.5), height=Inches(5.5))

    # Output
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


# =====================================================
# タブ構成
# =====================================================
tab_dashboard, tab_ai = st.tabs(["ダッシュボード", "AI問い合わせ"])

def _upload_and_get_url(data_bytes, file_name):
    """ステージにファイルをアップロードしPresigned URLを返す"""
    stage_path = f"@KFUKAMORI_GEN_DB.STADIA360.STADIA360_STAGE/reports/{file_name}"
    stream = io.BytesIO(data_bytes)
    session.file.put_stream(
        stream,
        stage_path,
        auto_compress=False,
        overwrite=True,
    )
    row = session.sql(
        f"SELECT GET_PRESIGNED_URL(@KFUKAMORI_GEN_DB.STADIA360.STADIA360_STAGE, 'reports/{file_name}', 3600) AS URL"
    ).collect()
    return row[0]["URL"]


with tab_dashboard:

    # --- レポート出力ボタン ---
    st.markdown("---")
    col_dl1, col_dl2, col_dl3, col_dl4 = st.columns([1, 1, 1, 3])
    today_str = datetime.date.today().strftime("%Y%m%d")
    with col_dl1:
        if st.button("PDF Export", key="pdf_gen_btn"):
            with st.spinner("PDF生成中..."):
                pdf_bytes = generate_pdf_report()
                pdf_name = f"STADIA360_Report_{today_str}.pdf"
                url = _upload_and_get_url(pdf_bytes, pdf_name)
                st.session_state["pdf_url"] = url
                st.session_state["pdf_name"] = pdf_name
        if "pdf_url" in st.session_state:
            st.markdown(f"[{st.session_state['pdf_name']} をダウンロード]({st.session_state['pdf_url']})")
    with col_dl2:
        if st.button("Excel Export", key="excel_gen_btn"):
            with st.spinner("Excel生成中..."):
                excel_bytes = generate_excel_report()
                excel_name = f"STADIA360_Report_{today_str}.xlsx"
                url = _upload_and_get_url(excel_bytes, excel_name)
                st.session_state["excel_url"] = url
                st.session_state["excel_name"] = excel_name
        if "excel_url" in st.session_state:
            st.markdown(f"[{st.session_state['excel_name']} をダウンロード]({st.session_state['excel_url']})")
    with col_dl3:
        if st.button("PPTX Export", key="pptx_gen_btn"):
            with st.spinner("PPTX生成中..."):
                pptx_bytes = generate_pptx_report()
                pptx_name = f"STADIA360_Report_{today_str}.pptx"
                url = _upload_and_get_url(pptx_bytes, pptx_name)
                st.session_state["pptx_url"] = url
                st.session_state["pptx_name"] = pptx_name
        if "pptx_url" in st.session_state:
            st.markdown(f"[{st.session_state['pptx_name']} をダウンロード]({st.session_state['pptx_url']})")
    st.markdown("---")

    # =====================================================
    # 1. 概況サマリー（KPI）
    # =====================================================
    st.header("概況サマリー")

    col1, col2, col3, col4 = st.columns(4)

    with col1:
        cm_count = df_tv_f[df_tv_f["CM_EXPOSED"] == True].shape[0]
        total_tv = df_tv_f.shape[0]
        cm_rate = (cm_count / total_tv * 100) if total_tv > 0 else 0
        st.metric("CM接触件数", f"{cm_count:,}", f"接触率 {cm_rate:.1f}%")

    with col2:
        cv_count = df_site_f[df_site_f["CONVERSION_FLAG"] == True].shape[0]
        total_sessions = df_site_f.shape[0]
        cvr = (cv_count / total_sessions * 100) if total_sessions > 0 else 0
        st.metric("サイトCV数", f"{cv_count:,}", f"CVR {cvr:.1f}%")

    with col3:
        store_visits = df_store_f.shape[0]
        avg_stay = df_store_f["STAY_MINUTES"].mean() if store_visits > 0 else 0
        st.metric("来店件数", f"{store_visits:,}", f"平均滞在 {avg_stay:.0f}分")

    with col4:
        dl_count = df_app_dl_f.shape[0]
        launch_count = df_app_launch_f.shape[0]
        st.metric("アプリDL数", f"{dl_count:,}", f"起動 {launch_count:,}回")

    col5, col6, col7, col8 = st.columns(4)

    with col5:
        total_purchase = df_purchase_f["AMOUNT"].sum()
        st.metric("購買総額", f"¥{total_purchase:,.0f}")

    with col6:
        avg_nps = df_loyalty["NPS_SCORE"].mean()
        st.metric("平均NPS", f"{avg_nps:.1f}")

    with col7:
        total_pv = df_site_f["PAGE_VIEWS"].sum()
        st.metric("合計PV", f"{total_pv:,}")

    with col8:
        total_hours = df_tv_f["VIEWING_SECONDS"].sum() / 3600
        st.metric("総視聴時間", f"{total_hours:,.0f}時間")


    # =====================================================
    # 2. テレビ視聴分析
    # =====================================================
    st.header("テレビ視聴分析")

    col_tv1, col_tv2 = st.columns(2)

    with col_tv1:
        st.subheader("チャンネル別視聴件数")
        ch_data = df_tv_f.groupby("CHANNEL").agg(
            VIEWS=("VIEWING_ID", "count"),
            CM_EXPOSED=("CM_EXPOSED", "sum")
        ).reset_index()
        ch_data["CM接触率(%)"] = (ch_data["CM_EXPOSED"] / ch_data["VIEWS"] * 100).round(1)

        chart_ch = (
            alt.Chart(ch_data)
            .mark_bar(color="#1f77b4")
            .encode(
                x=alt.X("VIEWS:Q", title="視聴件数"),
                y=alt.Y("CHANNEL:N", title="チャンネル", sort="-x"),
                tooltip=["CHANNEL", alt.Tooltip("VIEWS:Q", format=","), alt.Tooltip("CM接触率(%):Q", format=".1f")],
            )
        )
        st.altair_chart(chart_ch, use_container_width=True)

    with col_tv2:
        st.subheader("時間帯別視聴トレンド")
        hour_data = df_tv_f.groupby("VIEWING_HOUR").agg(
            VIEWS=("VIEWING_ID", "count")
        ).reset_index()

        chart_hour = (
            alt.Chart(hour_data)
            .mark_area(color="#ff7f0e", opacity=0.7, line=True)
            .encode(
                x=alt.X("VIEWING_HOUR:O", title="時間帯"),
                y=alt.Y("VIEWS:Q", title="視聴件数"),
                tooltip=["VIEWING_HOUR", alt.Tooltip("VIEWS:Q", format=",")],
            )
        )
        st.altair_chart(chart_hour, use_container_width=True)

    # TV vs CTV 比較
    st.subheader("TV vs CTV デバイス別視聴比較")
    col_dev1, col_dev2 = st.columns(2)

    with col_dev1:
        device_data = df_tv_f.groupby("DEVICE_TYPE").agg(
            VIEWS=("VIEWING_ID", "count")
        ).reset_index()
        chart_device = (
            alt.Chart(device_data)
            .mark_arc(innerRadius=50)
            .encode(
                theta=alt.Theta("VIEWS:Q"),
                color=alt.Color("DEVICE_TYPE:N", title="デバイス"),
                tooltip=["DEVICE_TYPE", alt.Tooltip("VIEWS:Q", format=",")],
            )
        )
        st.altair_chart(chart_device, use_container_width=True)

    with col_dev2:
        cm_by_device = df_tv_f.groupby("DEVICE_TYPE").agg(
            TOTAL=("VIEWING_ID", "count"),
            CM=("CM_EXPOSED", "sum")
        ).reset_index()
        cm_by_device["CM接触率(%)"] = (cm_by_device["CM"] / cm_by_device["TOTAL"] * 100).round(1)
        st.dataframe(cm_by_device.rename(columns={
            "DEVICE_TYPE": "デバイス", "TOTAL": "視聴件数", "CM": "CM接触件数"
        }), use_container_width=True)


    # =====================================================
    # CM効果分析（レスポンス分析）
    # =====================================================
    st.header("CM効果分析（レスポンス分析）")
    st.caption("CM接触者のサイトコンバージョン（CV）を軸にした効果検証")

    # CONVERSION_FLAG を数値型に変換（Snowflakeからobject型で返される場合がある）
    df_cm_cv_f["CONVERSION_FLAG"] = pd.to_numeric(df_cm_cv_f["CONVERSION_FLAG"], errors="coerce").fillna(0).astype(int)

    # --- 放送局別ドリル ---
    col_cm1, col_cm2 = st.columns(2)

    with col_cm1:
        st.subheader("放送局別ドリル（CV率比較）")
        ch_cm = df_cm_cv_f.groupby("CHANNEL").agg(
            CM接触数=("VIEWING_ID", "count"),
            CV数=("CONVERSION_FLAG", "sum"),
        ).reset_index()
        ch_cm["CV率(%)"] = (ch_cm["CV数"] / ch_cm["CM接触数"] * 100).round(2)

        chart_ch_cv = (
            alt.Chart(ch_cm)
            .mark_bar(color="#1f77b4")
            .encode(
                x=alt.X("CV率(%):Q", title="CV率（%）"),
                y=alt.Y("CHANNEL:N", title="放送局", sort="-x"),
                tooltip=["CHANNEL",
                          alt.Tooltip("CM接触数:Q", format=","),
                          alt.Tooltip("CV数:Q", format=","),
                          alt.Tooltip("CV率(%):Q", format=".2f")],
            )
        )
        st.altair_chart(chart_ch_cv, use_container_width=True)

    # --- クリエイティブドリル ---
    with col_cm2:
        st.subheader("クリエイティブドリル（素材別CV率）")
        cr_cm = df_cm_cv_f[df_cm_cv_f["CREATIVE_NAME"].notna()].groupby("CREATIVE_NAME").agg(
            CM接触数=("VIEWING_ID", "count"),
            CV数=("CONVERSION_FLAG", "sum"),
        ).reset_index()
        cr_cm["CV率(%)"] = (cr_cm["CV数"] / cr_cm["CM接触数"] * 100).round(2)

        chart_cr_cv = (
            alt.Chart(cr_cm)
            .mark_bar(color="#ff7f0e")
            .encode(
                x=alt.X("CV率(%):Q", title="CV率（%）"),
                y=alt.Y("CREATIVE_NAME:N", title="CM素材", sort="-x"),
                tooltip=["CREATIVE_NAME",
                          alt.Tooltip("CM接触数:Q", format=","),
                          alt.Tooltip("CV数:Q", format=","),
                          alt.Tooltip("CV率(%):Q", format=".2f")],
            )
        )
        st.altair_chart(chart_cr_cv, use_container_width=True)

    # --- レスポンスヒートマップ ---
    st.subheader("レスポンスヒートマップ（曜日×時間帯別CV数）")

    hm_channel = st.selectbox("放送局で絞り込み", ["全局"] + sorted(df_cm_cv_f["CHANNEL"].unique().tolist()), key="hm_ch")
    df_hm = df_cm_cv_f[df_cm_cv_f["CONVERSION_FLAG"] == True].copy()
    if hm_channel != "全局":
        df_hm = df_hm[df_hm["CHANNEL"] == hm_channel]

    dow_order = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    dow_label = {"Mon": "月", "Tue": "火", "Wed": "水", "Thu": "木", "Fri": "金", "Sat": "土", "Sun": "日"}
    df_hm["曜日"] = df_hm["CM_DOW"].map(dow_label)

    hm_data = df_hm.groupby(["曜日", "CM_HOUR"]).agg(
        CV数=("VIEWING_ID", "count")
    ).reset_index()
    hm_data.rename(columns={"CM_HOUR": "時間帯"}, inplace=True)

    dow_jp_order = ["月", "火", "水", "木", "金", "土", "日"]

    chart_hm = (
        alt.Chart(hm_data)
        .mark_rect()
        .encode(
            x=alt.X("曜日:N", title="曜日", sort=dow_jp_order),
            y=alt.Y("時間帯:O", title="時間帯", sort=list(range(24))),
            color=alt.Color("CV数:Q", title="CV数",
                            scale=alt.Scale(scheme="blues")),
            tooltip=["曜日", alt.Tooltip("時間帯:O"), alt.Tooltip("CV数:Q", format=",")],
        )
        .properties(height=400)
    )
    st.altair_chart(chart_hm, use_container_width=True)

    # --- レスポンス番組ランキング ---
    st.subheader("レスポンス番組ランキング")

    pgm_data = df_cm_cv_f.groupby("PROGRAM_NAME").agg(
        CM接触数=("VIEWING_ID", "count"),
        CV数=("CONVERSION_FLAG", "sum"),
    ).reset_index()
    pgm_data["CV率(%)"] = (pgm_data["CV数"] / pgm_data["CM接触数"] * 100).round(2)
    pgm_data = pgm_data.sort_values("CV率(%)", ascending=False).reset_index(drop=True)
    pgm_data.index = pgm_data.index + 1
    pgm_data.index.name = "順位"

    col_pgm1, col_pgm2 = st.columns([1, 1])

    with col_pgm1:
        chart_pgm = (
            alt.Chart(pgm_data.reset_index())
            .mark_bar(color="#2ca02c")
            .encode(
                x=alt.X("CV率(%):Q", title="CV率（%）"),
                y=alt.Y("PROGRAM_NAME:N", title="番組名", sort="-x"),
                tooltip=["PROGRAM_NAME",
                          alt.Tooltip("CM接触数:Q", format=","),
                          alt.Tooltip("CV数:Q", format=","),
                          alt.Tooltip("CV率(%):Q", format=".2f")],
            )
        )
        st.altair_chart(chart_pgm, use_container_width=True)

    with col_pgm2:
        st.dataframe(
            pgm_data.rename(columns={
                "PROGRAM_NAME": "番組名", "CM接触数": "CM接触数", "CV数": "CV数"
            })[["番組名", "CM接触数", "CV数", "CV率(%)"]],
            use_container_width=True,
        )

    # --- リーセンシー（CM接触→CV日数分布） ---
    st.subheader("リーセンシー（CM接触からCVまでの日数分布）")

    df_recency = df_cm_cv_f[
        (df_cm_cv_f["CONVERSION_FLAG"] == True) &
        (df_cm_cv_f["DAYS_TO_CV"].notna()) &
        (df_cm_cv_f["DAYS_TO_CV"] >= 0) &
        (df_cm_cv_f["DAYS_TO_CV"] <= 30)
    ].copy()

    recency_data = df_recency.groupby("DAYS_TO_CV").agg(
        CV件数=("VIEWING_ID", "count")
    ).reset_index()
    recency_data.rename(columns={"DAYS_TO_CV": "経過日数"}, inplace=True)

    chart_recency = (
        alt.Chart(recency_data)
        .mark_area(color="#9467bd", opacity=0.7, line=True)
        .encode(
            x=alt.X("経過日数:Q", title="CM接触からCVまでの経過日数",
                    scale=alt.Scale(domain=[0, 30])),
            y=alt.Y("CV件数:Q", title="CV件数"),
            tooltip=[alt.Tooltip("経過日数:Q"), alt.Tooltip("CV件数:Q", format=",")],
        )
        .properties(height=300)
    )
    st.altair_chart(chart_recency, use_container_width=True)


    # =====================================================
    # 3. 態度変容ファネル分析
    # =====================================================
    st.header("態度変容ファネル分析")

    col_att1, col_att2 = st.columns(2)

    with col_att1:
        st.subheader("CM接触者 vs 非接触者 リフト比較")
        exposed = df_attitude_f[df_attitude_f["CM_EXPOSED"] == True]
        unexposed = df_attitude_f[df_attitude_f["CM_EXPOSED"] == False]

        lift_data = pd.DataFrame({
            "ステージ": ["認知", "興味", "検討", "購入意向"],
            "CM接触者リフト": [
                (exposed["AWARENESS_AFTER"] - exposed["AWARENESS_BEFORE"]).mean() if len(exposed) > 0 else 0,
                (exposed["INTEREST_AFTER"] - exposed["INTEREST_BEFORE"]).mean() if len(exposed) > 0 else 0,
                (exposed["CONSIDER_AFTER"] - exposed["CONSIDER_BEFORE"]).mean() if len(exposed) > 0 else 0,
                (exposed["PURCHASE_AFTER"] - exposed["PURCHASE_BEFORE"]).mean() if len(exposed) > 0 else 0,
            ],
            "非接触者リフト": [
                (unexposed["AWARENESS_AFTER"] - unexposed["AWARENESS_BEFORE"]).mean() if len(unexposed) > 0 else 0,
                (unexposed["INTEREST_AFTER"] - unexposed["INTEREST_BEFORE"]).mean() if len(unexposed) > 0 else 0,
                (unexposed["CONSIDER_AFTER"] - unexposed["CONSIDER_BEFORE"]).mean() if len(unexposed) > 0 else 0,
                (unexposed["PURCHASE_AFTER"] - unexposed["PURCHASE_BEFORE"]).mean() if len(unexposed) > 0 else 0,
            ],
        })

        lift_long = lift_data.melt(id_vars=["ステージ"], var_name="グループ", value_name="リフト（pt）")
        stage_order = ["認知", "興味", "検討", "購入意向"]

        chart_lift = (
            alt.Chart(lift_long)
            .mark_bar()
            .encode(
                x=alt.X("グループ:N", title=""),
                y=alt.Y("リフト（pt）:Q", title="平均リフト（ポイント）"),
                color=alt.Color("グループ:N", scale=alt.Scale(
                    domain=["CM接触者リフト", "非接触者リフト"],
                    range=["#2ca02c", "#d62728"]
                )),
                column=alt.Column("ステージ:N", title="ステージ", sort=stage_order),
                tooltip=["ステージ", "グループ", alt.Tooltip("リフト（pt）:Q", format=".2f")],
            )
            .properties(width=100)
        )
        st.altair_chart(chart_lift, use_container_width=True)

    with col_att2:
        st.subheader("キャンペーン別 認知リフト")
        att_by_camp = df_attitude_f.merge(
            df_campaigns[["CAMPAIGN_ID", "CAMPAIGN_NAME"]], on="CAMPAIGN_ID", how="left"
        )
        camp_lift = att_by_camp.groupby("CAMPAIGN_NAME").apply(
            lambda x: pd.Series({
                "認知リフト": (x["AWARENESS_AFTER"] - x["AWARENESS_BEFORE"]).mean(),
                "購入リフト": (x["PURCHASE_AFTER"] - x["PURCHASE_BEFORE"]).mean(),
                "回答数": len(x),
            })
        ).reset_index()

        chart_camp_lift = (
            alt.Chart(camp_lift)
            .mark_bar(color="#1f77b4")
            .encode(
                x=alt.X("認知リフト:Q", title="平均認知リフト（pt）"),
                y=alt.Y("CAMPAIGN_NAME:N", title="キャンペーン", sort="-x"),
                tooltip=["CAMPAIGN_NAME", alt.Tooltip("認知リフト:Q", format=".2f"),
                          alt.Tooltip("購入リフト:Q", format=".2f"), alt.Tooltip("回答数:Q", format=",")],
            )
        )
        st.altair_chart(chart_camp_lift, use_container_width=True)


    # =====================================================
    # 4. サイト来訪分析
    # =====================================================
    st.header("サイト来訪分析")

    col_s1, col_s2 = st.columns(2)

    with col_s1:
        st.subheader("流入経路別セッション数・CVR")
        ref_data = df_site_f.groupby("REFERRER_TYPE").agg(
            SESSIONS=("SESSION_ID", "count"),
            CVS=("CONVERSION_FLAG", "sum"),
            AVG_PV=("PAGE_VIEWS", "mean"),
        ).reset_index()
        ref_data["CVR(%)"] = (ref_data["CVS"] / ref_data["SESSIONS"] * 100).round(1)
        ref_data["AVG_PV"] = ref_data["AVG_PV"].round(1)

        chart_ref = (
            alt.Chart(ref_data)
            .mark_bar(color="#1f77b4")
            .encode(
                x=alt.X("SESSIONS:Q", title="セッション数"),
                y=alt.Y("REFERRER_TYPE:N", title="流入経路", sort="-x"),
                tooltip=["REFERRER_TYPE", alt.Tooltip("SESSIONS:Q", format=","),
                          alt.Tooltip("CVR(%):Q", format=".1f"), alt.Tooltip("AVG_PV:Q", format=".1f")],
            )
        )
        st.altair_chart(chart_ref, use_container_width=True)

    with col_s2:
        st.subheader("月別セッション数・CV推移")
        df_site_f_copy = df_site_f.copy()
        df_site_f_copy["MONTH"] = pd.to_datetime(df_site_f_copy["VISIT_DATE"]).dt.to_period("M").astype(str)
        monthly_site = df_site_f_copy.groupby("MONTH").agg(
            SESSIONS=("SESSION_ID", "count"),
            CVS=("CONVERSION_FLAG", "sum"),
        ).reset_index()

        base = alt.Chart(monthly_site).encode(x=alt.X("MONTH:N", title="月"))
        bar = base.mark_bar(color="#1f77b4", opacity=0.6).encode(
            y=alt.Y("SESSIONS:Q", title="セッション数"),
            tooltip=["MONTH", alt.Tooltip("SESSIONS:Q", format=",")],
        )
        line = base.mark_line(color="#d62728", strokeWidth=2, point=True).encode(
            y=alt.Y("CVS:Q", title="CV数"),
            tooltip=["MONTH", alt.Tooltip("CVS:Q", format=",")],
        )
        chart_monthly = alt.layer(bar, line).resolve_scale(y="independent")
        st.altair_chart(chart_monthly, use_container_width=True)


    # =====================================================
    # 5. オフライン購買分析
    # =====================================================
    st.header("オフライン購買分析")

    col_p1, col_p2 = st.columns(2)

    with col_p1:
        st.subheader("商品カテゴリ別購買額")
        cat_data = df_purchase_f.groupby("PRODUCT_CATEGORY").agg(
            TOTAL_AMOUNT=("AMOUNT", "sum"),
            COUNT=("PURCHASE_ID", "count"),
        ).reset_index()
        cat_data["AVG_AMOUNT"] = (cat_data["TOTAL_AMOUNT"] / cat_data["COUNT"]).round(0)

        chart_cat = (
            alt.Chart(cat_data)
            .mark_bar(color="#2ca02c")
            .encode(
                x=alt.X("TOTAL_AMOUNT:Q", title="購買総額（円）"),
                y=alt.Y("PRODUCT_CATEGORY:N", title="カテゴリ", sort="-x"),
                tooltip=["PRODUCT_CATEGORY", alt.Tooltip("TOTAL_AMOUNT:Q", format=","),
                          alt.Tooltip("COUNT:Q", format=","), alt.Tooltip("AVG_AMOUNT:Q", format=",")],
            )
        )
        st.altair_chart(chart_cat, use_container_width=True)

    with col_p2:
        st.subheader("CM接触者 vs 非接触者 購買比較")
        cm_purchase = df_purchase_f.groupby("CM_EXPOSED").agg(
            AVG_AMOUNT=("AMOUNT", "mean"),
            COUNT=("PURCHASE_ID", "count"),
            TOTAL=("AMOUNT", "sum"),
        ).reset_index()
        cm_purchase["グループ"] = cm_purchase["CM_EXPOSED"].map({True: "CM接触者", False: "非接触者"})
        cm_purchase["AVG_AMOUNT"] = cm_purchase["AVG_AMOUNT"].round(0)

        chart_cm_purchase = (
            alt.Chart(cm_purchase)
            .mark_bar()
            .encode(
                x=alt.X("グループ:N", title=""),
                y=alt.Y("AVG_AMOUNT:Q", title="平均購買金額（円）"),
                color=alt.Color("グループ:N", scale=alt.Scale(
                    domain=["CM接触者", "非接触者"],
                    range=["#2ca02c", "#d62728"]
                )),
                tooltip=["グループ", alt.Tooltip("AVG_AMOUNT:Q", format=","),
                          alt.Tooltip("COUNT:Q", format=","), alt.Tooltip("TOTAL:Q", format=",")],
            )
        )
        st.altair_chart(chart_cm_purchase, use_container_width=True)


    # =====================================================
    # 6. 来店・アプリ分析
    # =====================================================
    st.header("来店・アプリ分析")

    col_a1, col_a2 = st.columns(2)

    with col_a1:
        st.subheader("店舗別来店件数")
        store_data = df_store_f.groupby("STORE_NAME").agg(
            VISITS=("VISIT_ID", "count"),
            AVG_STAY=("STAY_MINUTES", "mean"),
        ).reset_index()
        store_data["AVG_STAY"] = store_data["AVG_STAY"].round(1)

        chart_store = (
            alt.Chart(store_data)
            .mark_bar(color="#9467bd")
            .encode(
                x=alt.X("VISITS:Q", title="来店件数"),
                y=alt.Y("STORE_NAME:N", title="店舗", sort="-x"),
                tooltip=["STORE_NAME", alt.Tooltip("VISITS:Q", format=","),
                          alt.Tooltip("AVG_STAY:Q", format=".1f")],
            )
        )
        st.altair_chart(chart_store, use_container_width=True)

    with col_a2:
        st.subheader("アプリ別DL数・起動数")
        dl_by_app = df_app_dl_f.groupby("APP_NAME").agg(DL=("DOWNLOAD_ID", "count")).reset_index()
        launch_by_app = df_app_launch_f.groupby("APP_NAME").agg(LAUNCH=("LAUNCH_ID", "count")).reset_index()
        app_data = dl_by_app.merge(launch_by_app, on="APP_NAME", how="outer").fillna(0)
        app_data["LAUNCH"] = app_data["LAUNCH"].astype(int)
        app_data["DL"] = app_data["DL"].astype(int)

        app_long = app_data.melt(id_vars=["APP_NAME"], var_name="指標", value_name="件数")

        chart_app = (
            alt.Chart(app_long)
            .mark_bar()
            .encode(
                x=alt.X("指標:N", title=""),
                y=alt.Y("件数:Q", title="件数"),
                color=alt.Color("指標:N", scale=alt.Scale(
                    domain=["DL", "LAUNCH"],
                    range=["#1f77b4", "#ff7f0e"]
                )),
                column=alt.Column("APP_NAME:N", title="アプリ"),
                tooltip=["APP_NAME", "指標", alt.Tooltip("件数:Q", format=",")],
            )
            .properties(width=100)
        )
        st.altair_chart(chart_app, use_container_width=True)

    # アプリDLチャネル分析
    st.subheader("アプリDL 広告チャネル別")
    col_ch1, col_ch2 = st.columns(2)

    with col_ch1:
        channel_data = df_app_dl_f.groupby("AD_CHANNEL").agg(
            DL=("DOWNLOAD_ID", "count")
        ).reset_index()
        chart_channel = (
            alt.Chart(channel_data)
            .mark_arc(innerRadius=50)
            .encode(
                theta=alt.Theta("DL:Q"),
                color=alt.Color("AD_CHANNEL:N", title="チャネル"),
                tooltip=["AD_CHANNEL", alt.Tooltip("DL:Q", format=",")],
            )
        )
        st.altair_chart(chart_channel, use_container_width=True)

    with col_ch2:
        os_data = df_app_dl_f.groupby("OS_TYPE").agg(DL=("DOWNLOAD_ID", "count")).reset_index()
        chart_os = (
            alt.Chart(os_data)
            .mark_arc(innerRadius=50)
            .encode(
                theta=alt.Theta("DL:Q"),
                color=alt.Color("OS_TYPE:N", title="OS",
                                scale=alt.Scale(domain=["iOS", "Android"], range=["#636363", "#2ca02c"])),
                tooltip=["OS_TYPE", alt.Tooltip("DL:Q", format=",")],
            )
        )
        st.altair_chart(chart_os, use_container_width=True)


    # =====================================================
    # 7. 顧客ロイヤリティ分析
    # =====================================================
    st.header("顧客ロイヤリティ分析")

    col_l1, col_l2 = st.columns(2)

    with col_l1:
        st.subheader("ロイヤリティセグメント分布")
        seg_data = df_loyalty.groupby("LOYALTY_SEGMENT").agg(
            COUNT=("CUSTOMER_ID", "count"),
            AVG_NPS=("NPS_SCORE", "mean"),
            AVG_LTV=("LTV_AMOUNT", "mean"),
        ).reset_index()
        seg_data["AVG_NPS"] = seg_data["AVG_NPS"].round(1)
        seg_data["AVG_LTV"] = seg_data["AVG_LTV"].round(0)

        chart_seg = (
            alt.Chart(seg_data)
            .mark_bar()
            .encode(
                x=alt.X("LOYALTY_SEGMENT:N", title="セグメント",
                         sort=["プロモーター", "パッシブ", "デトラクター"]),
                y=alt.Y("COUNT:Q", title="顧客数"),
                color=alt.Color("LOYALTY_SEGMENT:N", scale=alt.Scale(
                    domain=["プロモーター", "パッシブ", "デトラクター"],
                    range=["#2ca02c", "#ff7f0e", "#d62728"]
                ), legend=None),
                tooltip=["LOYALTY_SEGMENT", alt.Tooltip("COUNT:Q", format=","),
                          alt.Tooltip("AVG_NPS:Q", format=".1f"), alt.Tooltip("AVG_LTV:Q", format=",")],
            )
        )
        st.altair_chart(chart_seg, use_container_width=True)

    with col_l2:
        st.subheader("年齢層 × 性別 分布")
        age_gender = df_loyalty.groupby(["AGE_GROUP", "GENDER"]).agg(
            COUNT=("CUSTOMER_ID", "count")
        ).reset_index()
        age_order = ["10代", "20代", "30代", "40代", "50代", "60代"]

        chart_age = (
            alt.Chart(age_gender)
            .mark_bar()
            .encode(
                x=alt.X("GENDER:N", title=""),
                y=alt.Y("COUNT:Q", title="顧客数"),
                color=alt.Color("GENDER:N", title="性別", scale=alt.Scale(
                    domain=["男性", "女性"],
                    range=["#1f77b4", "#e377c2"]
                )),
                column=alt.Column("AGE_GROUP:N", title="年齢層", sort=age_order),
                tooltip=["AGE_GROUP", "GENDER", alt.Tooltip("COUNT:Q", format=",")],
            )
            .properties(width=80)
        )
        st.altair_chart(chart_age, use_container_width=True)

    # エリア別NPS
    st.subheader("エリア別 平均NPS・平均LTV")
    area_loyalty = df_loyalty[df_loyalty["AREA"].isin(selected_areas)].groupby("AREA").agg(
        AVG_NPS=("NPS_SCORE", "mean"),
        AVG_LTV=("LTV_AMOUNT", "mean"),
        COUNT=("CUSTOMER_ID", "count"),
    ).reset_index()
    area_loyalty["AVG_NPS"] = area_loyalty["AVG_NPS"].round(1)
    area_loyalty["AVG_LTV"] = area_loyalty["AVG_LTV"].round(0)

    st.dataframe(
        area_loyalty.rename(columns={
            "AREA": "エリア", "AVG_NPS": "平均NPS", "AVG_LTV": "平均LTV（円）", "COUNT": "顧客数"
        }),
        use_container_width=True,
    )


# =====================================================
# AI問い合わせタブ
# =====================================================
with tab_ai:
    st.header("AI問い合わせ")
    st.caption("STADIA360のデータに基づいて、自然言語で質問できます（Snowflake Cortex COMPLETE）")

    # --- データコンテキスト生成 ---
    def build_data_context():
        """フィルタ適用済みの各テーブルから主要指標をサマリーテキスト化"""
        lines = []
        lines.append("=== STADIA360 統合マーケティングデータサマリー ===")

        # フィルター状態
        lines.append(f"\n【フィルター条件】")
        lines.append(f"- キャンペーン: {selected_campaign}")
        lines.append(f"- エリア: {', '.join(selected_areas)}")

        # キャンペーン一覧
        lines.append(f"\n【キャンペーン一覧】({len(df_campaigns)}件)")
        for _, r in df_campaigns.iterrows():
            lines.append(f"- {r['CAMPAIGN_NAME']} (ID:{r['CAMPAIGN_ID']}, {r['START_DATE']}~{r['END_DATE']}, 予算:{r['BUDGET_MM']}百万円)")

        # TV視聴
        lines.append(f"\n【TV視聴データ】(フィルタ後 {len(df_tv_f):,}件)")
        cm_cnt = df_tv_f[df_tv_f["CM_EXPOSED"] == True].shape[0]
        lines.append(f"- CM接触件数: {cm_cnt:,} / 接触率: {(cm_cnt/len(df_tv_f)*100 if len(df_tv_f)>0 else 0):.1f}%")
        ch_top = df_tv_f.groupby("CHANNEL").size().sort_values(ascending=False).head(5)
        lines.append(f"- チャンネル別視聴TOP5: {', '.join(f'{k}({v:,}件)' for k,v in ch_top.items())}")
        dev_dist = df_tv_f.groupby("DEVICE_TYPE").size()
        lines.append(f"- デバイス別: {', '.join(f'{k}({v:,}件)' for k,v in dev_dist.items())}")

        # CM効果(CM→CV)
        if len(df_cm_cv_f) > 0:
            cv_flag = pd.to_numeric(df_cm_cv_f["CONVERSION_FLAG"], errors="coerce").fillna(0).astype(int)
            cv_total = cv_flag.sum()
            cv_rate = (cv_total / len(df_cm_cv_f) * 100)
            lines.append(f"\n【CM効果分析】(CM接触者 {len(df_cm_cv_f):,}件)")
            lines.append(f"- CM接触者CV数: {cv_total:,} / CV率: {cv_rate:.2f}%")
            ch_cv = df_cm_cv_f.copy()
            ch_cv["CV"] = cv_flag
            ch_grp = ch_cv.groupby("CHANNEL").agg(接触数=("VIEWING_ID","count"), CV数=("CV","sum")).reset_index()
            ch_grp["CV率"] = (ch_grp["CV数"]/ch_grp["接触数"]*100).round(2)
            for _, r in ch_grp.sort_values("CV率", ascending=False).iterrows():
                lines.append(f"  - {r['CHANNEL']}: 接触{r['接触数']:,}件, CV{r['CV数']:,}件, CV率{r['CV率']:.2f}%")

        # 態度変容
        if len(df_attitude_f) > 0:
            exp = df_attitude_f[df_attitude_f["CM_EXPOSED"] == True]
            unexp = df_attitude_f[df_attitude_f["CM_EXPOSED"] == False]
            lines.append(f"\n【態度変容】(フィルタ後 {len(df_attitude_f):,}件)")
            if len(exp) > 0:
                aw_lift = (exp["AWARENESS_AFTER"] - exp["AWARENESS_BEFORE"]).mean()
                pu_lift = (exp["PURCHASE_AFTER"] - exp["PURCHASE_BEFORE"]).mean()
                lines.append(f"- CM接触者: 認知リフト{aw_lift:.2f}pt, 購入意向リフト{pu_lift:.2f}pt")
            if len(unexp) > 0:
                aw_lift_u = (unexp["AWARENESS_AFTER"] - unexp["AWARENESS_BEFORE"]).mean()
                pu_lift_u = (unexp["PURCHASE_AFTER"] - unexp["PURCHASE_BEFORE"]).mean()
                lines.append(f"- 非接触者: 認知リフト{aw_lift_u:.2f}pt, 購入意向リフト{pu_lift_u:.2f}pt")

        # サイト来訪
        if len(df_site_f) > 0:
            site_cv = df_site_f[df_site_f["CONVERSION_FLAG"] == True].shape[0]
            site_cvr = (site_cv / len(df_site_f) * 100)
            lines.append(f"\n【サイト来訪】(フィルタ後 {len(df_site_f):,}件)")
            lines.append(f"- CV数: {site_cv:,} / CVR: {site_cvr:.1f}%")
            lines.append(f"- 合計PV: {df_site_f['PAGE_VIEWS'].sum():,}")
            ref_top = df_site_f.groupby("REFERRER_TYPE").size().sort_values(ascending=False)
            lines.append(f"- 流入経路: {', '.join(f'{k}({v:,}件)' for k,v in ref_top.items())}")

        # オフライン購買
        if len(df_purchase_f) > 0:
            lines.append(f"\n【オフライン購買】(フィルタ後 {len(df_purchase_f):,}件)")
            lines.append(f"- 購買総額: ¥{df_purchase_f['AMOUNT'].sum():,.0f}")
            lines.append(f"- 平均購買額: ¥{df_purchase_f['AMOUNT'].mean():,.0f}")
            cat_top = df_purchase_f.groupby("PRODUCT_CATEGORY")["AMOUNT"].sum().sort_values(ascending=False)
            lines.append(f"- カテゴリ別: {', '.join(f'{k}(¥{v:,.0f})' for k,v in cat_top.items())}")
            cm_exp_avg = df_purchase_f[df_purchase_f["CM_EXPOSED"]==True]["AMOUNT"].mean()
            cm_unexp_avg = df_purchase_f[df_purchase_f["CM_EXPOSED"]==False]["AMOUNT"].mean()
            lines.append(f"- CM接触者平均: ¥{cm_exp_avg:,.0f} / 非接触者平均: ¥{cm_unexp_avg:,.0f}")

        # 来店
        if len(df_store_f) > 0:
            lines.append(f"\n【来店】(フィルタ後 {len(df_store_f):,}件)")
            lines.append(f"- 平均滞在: {df_store_f['STAY_MINUTES'].mean():.0f}分")

        # アプリ
        lines.append(f"\n【アプリ】DL:{len(df_app_dl_f):,}件 / 起動:{len(df_app_launch_f):,}件")

        # 顧客ロイヤリティ
        lines.append(f"\n【顧客ロイヤリティ】({len(df_loyalty):,}人)")
        lines.append(f"- 平均NPS: {df_loyalty['NPS_SCORE'].mean():.1f}")
        lines.append(f"- 平均LTV: ¥{df_loyalty['LTV_AMOUNT'].mean():,.0f}")
        seg_dist = df_loyalty.groupby("LOYALTY_SEGMENT").size()
        lines.append(f"- セグメント: {', '.join(f'{k}({v:,}人)' for k,v in seg_dist.items())}")

        return "\n".join(lines)

    # --- チャット履歴管理 ---
    if "messages" not in st.session_state:
        st.session_state.messages = []

    # --- サンプル質問 ---
    st.markdown("**質問例:**")
    sample_cols = st.columns(3)
    sample_questions = [
        "CM接触者と非接触者の購買金額の差を教えて",
        "最もCV率が高い放送局とその理由を分析して",
        "キャンペーン全体の効果をサマリーして",
    ]
    for i, q in enumerate(sample_questions):
        with sample_cols[i]:
            if st.button(q, key=f"sample_{i}"):
                st.session_state["ai_query_input"] = q

    st.markdown("---")

    # --- ユーザー入力 ---
    col_input, col_btn = st.columns([5, 1])
    with col_input:
        user_input = st.text_input(
            "質問を入力",
            value=st.session_state.get("ai_query_input", ""),
            key="ai_text_input",
            label_visibility="collapsed",
            placeholder="データについて質問してください...",
        )
    with col_btn:
        send_clicked = st.button("送信", key="ai_send_btn", type="primary")

    # --- 送信処理 ---
    if send_clicked and user_input:
        st.session_state.messages.append({"role": "user", "content": user_input})
        st.session_state["ai_query_input"] = ""

        with st.spinner("分析中..."):
            data_context = build_data_context()

            system_prompt = f"""あなたはSTADIA360統合マーケティングプラットフォームのデータアナリストです。
以下のデータサマリーに基づいて、ユーザーの質問に日本語で丁寧に回答してください。
具体的な数値を引用し、マーケティング施策への示唆を含めてください。

{data_context}"""

            escaped_system = system_prompt.replace("'", "''").replace("\\", "\\\\")
            escaped_question = user_input.replace("'", "''").replace("\\", "\\\\")

            prompt_for_complete = f"{escaped_system}\n\nユーザーの質問: {escaped_question}"

            try:
                result = session.sql(
                    f"SELECT SNOWFLAKE.CORTEX.COMPLETE('claude-3-5-sonnet', '{prompt_for_complete}') AS RESPONSE"
                ).collect()
                response_text = result[0]["RESPONSE"]
            except Exception as e:
                response_text = f"エラーが発生しました: {str(e)}"

            st.session_state.messages.append({"role": "assistant", "content": response_text})

    # --- チャット履歴表示 ---
    for msg in reversed(st.session_state.messages):
        if msg["role"] == "user":
            st.markdown(f"**Q:** {msg['content']}")
        else:
            st.markdown(f"**A:** {msg['content']}")
        st.markdown("---")
